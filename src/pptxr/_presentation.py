"""Presentation wrapper for python-pptx."""

from pathlib import Path
from typing import IO, Union, cast

from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.shapes.autoshape import Shape as PptxShape
from pptx.shapes.base import BaseShape
from pptx.slide import Slide as PptxSlide
from pptx.util import Pt

from ._data import Presentation, Shape, Slide
from .types import LiteralLength, to_points


class PptxShapeWrapper:
    """Wrapper for python-pptx Shape."""

    def __init__(self, shape: BaseShape) -> None:
        """Initialize the wrapper."""
        self._shape = cast(PptxShape, shape)

    def set_text(self, text: str) -> None:
        """Set the text of the shape."""
        self._shape.text = text


class PptxSlideWrapper:
    """Wrapper for python-pptx Slide."""

    def __init__(self, slide: PptxSlide) -> None:
        """Initialize the wrapper."""
        self._slide = slide

    def get_shapes(self) -> list[PptxShapeWrapper]:
        """Get all shapes in the slide."""
        return [PptxShapeWrapper(shape) for shape in self._slide.shapes]

    def add_shape(
        self,
        type: MSO_AUTO_SHAPE_TYPE,
        left: LiteralLength,
        top: LiteralLength,
        width: LiteralLength,
        height: LiteralLength,
    ) -> PptxShapeWrapper:
        """Add a shape to the slide."""
        shape = self._slide.shapes.add_shape(
            type,
            Pt(to_points(left)),
            Pt(to_points(top)),
            Pt(to_points(width)),
            Pt(to_points(height)),
        )
        return PptxShapeWrapper(shape)


class PptxPresentationWrapper:
    """Wrapper for python-pptx Presentation."""

    def __init__(self) -> None:
        """Initialize the wrapper."""
        self._presentation = PptxPresentation()

    def add_slide(self, slide: Slide) -> None:
        """Add a slide to the presentation."""
        pptx_slide = self._presentation.slides.add_slide(
            self._presentation.slide_layouts[0]
        )
        for shape in slide.shapes:
            self._add_shape(pptx_slide, shape)

    def _add_shape(self, pptx_slide: PptxSlide, shape: Shape) -> None:
        """Add a shape to a slide."""
        pptx_shape = pptx_slide.shapes.add_shape(
            shape.type,
            Pt(to_points(shape.left)),
            Pt(to_points(shape.top)),
            Pt(to_points(shape.width)),
            Pt(to_points(shape.height)),
        )
        if shape.text is not None:
            pptx_shape.text = shape.text

    def save(self, path: Union[str, Path, IO[bytes]]) -> None:
        """Save the presentation."""
        if isinstance(path, (str, Path)):
            self._presentation.save(str(path))
        else:
            self._presentation.save(path)


def create_presentation() -> "Presentation":
    """Create a new presentation.

    Returns:
        A new presentation.
    """
    wrapper = PptxPresentationWrapper()
    presentation = Presentation()
    presentation._wrapper = wrapper
    return presentation


def open_presentation(file: Union[str, Path]) -> "Presentation":
    """Open a presentation from a file.

    Args:
        file: The file to open.

    Returns:
        The opened presentation.
    """
    if isinstance(file, Path):
        file = str(file)
    wrapper = PptxPresentationWrapper()
    presentation = Presentation()
    presentation._wrapper = wrapper
    return presentation
