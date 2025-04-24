import os
import pathlib
from dataclasses import dataclass
from enum import Enum
from typing import (
    IO,
    Any,
    Literal,
    NotRequired,
    TypedDict,
    Unpack,
    overload,
)

import pptx.util
from pptx import Presentation as PptxPresentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN

from .units import (
    Length,
    _Inch,
    _to_internal_length,
    to_inche,
    to_point,
)


class SlideLayout(Enum):
    """Enumeration defining slide layout types"""

    TITLE = 0
    """Title-only layout"""

    TITLE_AND_CONTENT = 1
    """Title and content layout"""

    SECTION_HEADER = 2
    """Section header layout"""

    TWO_CONTENT = 3
    """Layout with two content areas"""

    COMPARISON = 4
    """Comparison layout"""

    TITLE_ONLY = 5
    """Title-only layout"""

    BLANK = 6
    """Blank layout"""

    CONTENT_WITH_CAPTION = 7
    """Content with caption layout"""

    PICTURE_WITH_CAPTION = 8
    """Picture with caption layout"""

    TITLE_AND_VERTICAL_TEXT = 9
    """Title and vertical text layout"""

    VERTICAL_TITLE_AND_TEXT = 10
    """Vertical title and text layout"""


class LayoutType(Enum):
    """Enumeration defining layout types"""

    FLEX = "flex"
    """Flexbox layout"""

    GRID = "grid"
    """Grid layout"""

    ABSOLUTE = "absolute"
    """Absolute positioning layout"""


class Align(Enum):
    """Enumeration defining element alignment"""

    START = "start"
    """Align to start position"""

    CENTER = "center"
    """Align to center"""

    END = "end"
    """Align to end position"""


class Justify(Enum):
    """Enumeration defining element justification"""

    START = "start"
    """Justify to start"""

    CENTER = "center"
    """Justify to center"""

    END = "end"
    """Justify to end"""

    SPACE_BETWEEN = "space-between"
    """Distribute space between elements"""

    SPACE_AROUND = "space-around"
    """Distribute space around elements"""


class Layout(TypedDict):
    """Data class representing layout settings"""

    type: NotRequired[LayoutType]
    """Layout type"""

    direction: NotRequired[str]
    """Layout direction ("row" or "column")"""

    align: NotRequired[Align]
    """Element alignment"""

    justify: NotRequired[Justify]
    """Element justification"""

    gap: NotRequired[Length]
    """Gap between elements"""

    padding: NotRequired[dict[str, Length]]
    """Padding (top, right, bottom, left)"""

    width: NotRequired[Length]
    """Width"""

    height: NotRequired[Length]
    """Height"""


class Text(TypedDict):
    """Data class representing text element"""

    type: Literal["text"]
    """Type of component"""

    text: str | None
    """Text content"""

    size: Length | None
    """Font size"""

    bold: NotRequired[bool]
    """Whether text is bold"""

    italic: NotRequired[bool]
    """Whether text is italic"""

    color: str | None
    """Text color"""

    layout: NotRequired[Layout]
    """Layout settings"""


@dataclass
class Shape:
    """Data class representing shape"""

    type: str
    """Shape type"""

    left: Length
    """Position from left edge"""

    top: Length
    """Position from top edge"""

    width: Length
    """Width"""

    height: Length
    """Height"""

    text: Text | None = None
    """Text within shape"""


class Image(TypedDict):
    """Data class representing image element"""

    type: Literal["image"]
    """Type of component"""

    path: str
    """Path to image file"""

    width: NotRequired[Length]
    """Width"""

    height: NotRequired[Length]
    """Height"""

    layout: NotRequired[Layout]
    """Layout settings"""


class Chart(TypedDict):
    """Data class representing chart element"""

    type: Literal["chart"]
    """Type of component"""

    chart_type: str
    """Chart type ("bar", "line", "pie", etc.)"""

    data: list[dict[str, Any]]
    """Chart data"""

    width: NotRequired[Length]
    """Width"""

    height: NotRequired[Length]
    """Height"""

    layout: NotRequired[Layout]
    """Layout settings"""


@dataclass
class TableCell:
    """Data class representing table cell"""

    text: str
    """Cell text content"""

    size: Length | None = None
    """Font size"""

    bold: bool = False
    """Whether text is bold"""

    italic: bool = False
    """Whether text is italic"""

    color: str | None = None
    """Text color"""

    background: str | None = None
    """Background color"""

    align: PP_ALIGN = PP_ALIGN.LEFT
    """Text alignment"""


class Table(TypedDict):
    """Data class representing table element"""

    type: Literal["table"]
    """Type of component"""

    rows: int
    """Number of rows"""

    cols: int
    """Number of columns"""

    data: list[list[TableCell]]
    """Table data"""

    width: NotRequired[Length]
    """Width"""

    height: NotRequired[Length]
    """Height"""

    layout: NotRequired[Layout]
    """Layout settings"""


Component = Text | Image | Chart | Table
"""Union type representing any component type"""


class Container(TypedDict):
    """Data class representing container"""

    components: list[Component]
    """Components within container"""

    layout: Layout
    """Layout settings"""


class Slide(TypedDict):
    """Data class representing slide"""

    layout: SlideLayout
    """Slide layout"""

    title: NotRequired[Text]
    """Slide title"""

    containers: NotRequired[list[Container]]
    """Containers within slide"""


class SlideTemplate:
    """Interface for slide templates

    This interface defines the contract for slide templates.
    Users can implement their own templates by subclassing this class
    and implementing the build method.
    """

    def build(self, title: Text | None = None, **kwargs) -> Slide:
        """Build a slide using this template

        Args:
            title (Text | None): Slide title
            **kwargs: Template-specific parameters

        Returns:
            Slide: Built slide
        """
        raise NotImplementedError()


class Presentation:
    """Class for creating presentations"""

    def __init__(self):
        """Initialize presentation"""
        self._presentation = PptxPresentation()

    @property
    def slide_layouts(self):
        """Get slide layouts"""
        return self._presentation.slide_layouts

    @property
    def slides(self):
        """Get slides"""
        return self._presentation.slides

    def save(self, path: str | pathlib.Path | IO[bytes]) -> None:
        """Save presentation to file

        Args:
            path (str): Path to save presentation
        """
        if isinstance(path, os.PathLike):
            path = str(path)

        self._presentation.save(path)

    @classmethod
    def builder(cls) -> "_PresentationBuilder":
        """Create a new presentation builder

        Returns:
            PresentationBuilder: Newly created presentation builder
        """
        return _PresentationBuilder()


class _PresentationBuilder:
    """Builder class for creating presentations"""

    def __init__(self):
        """Initialize presentation builder"""
        self.presentation = Presentation()
        self.slides: list[Slide] = []

    @overload
    def add_slide(self, slide: Slide | SlideTemplate, /) -> "_PresentationBuilder": ...

    @overload
    def add_slide(self, /, **kwargs: Unpack[Slide]) -> "_PresentationBuilder": ...

    def add_slide(  # type: ignore
        self,
        slide: Slide | SlideTemplate | None = None,
        /,
        **kwargs: Unpack[Slide],
    ) -> "_PresentationBuilder":
        """Add a slide to the presentation"""
        if isinstance(slide, SlideTemplate):
            # Using template
            slide = slide.build(**kwargs)
        elif slide is None:
            slide = kwargs

        self.slides.append(slide)
        return self

    def _apply_layout(self, shape, layout: Layout):
        """Apply layout to a shape

        Args:
            shape: Target shape
            layout (Layout): Layout settings to apply
        """
        # レイアウト情報を一度だけ取得
        width = layout.get("width")
        height = layout.get("height")
        align = layout.get("align")

        if width:
            internal_width = _to_internal_length(width)
            shape.width = pptx.util.Inches(to_inche(internal_width).value)
        if height:
            internal_height = _to_internal_length(height)
            shape.height = pptx.util.Inches(to_inche(internal_height).value)
        if align:
            if align == Align.CENTER:
                internal_left = _to_internal_length(shape.left)
                internal_width = _to_internal_length(shape.width)
                shape.left = (
                    pptx.util.Inches(to_inche(internal_left).value)
                    + (
                        pptx.util.Inches(8.5)
                        - pptx.util.Inches(to_inche(internal_width).value)
                    )
                    / 2
                )
            elif align == Align.END:
                internal_left = _to_internal_length(shape.left)
                internal_width = _to_internal_length(shape.width)
                shape.left = (
                    pptx.util.Inches(to_inche(internal_left).value)
                    + pptx.util.Inches(8.5)
                    - pptx.util.Inches(to_inche(internal_width).value)
                )

    def _add_component(
        self, slide_obj, component: Component, left: Length, top: Length
    ):
        """Add a component to a slide

        Args:
            slide_obj: Target slide
            component (Component): Component to add
            left (Length): Position from left edge
            top (Length): Position from top edge
        """
        # 共通の位置計算を最適化
        internal_left = _to_internal_length(left)
        internal_top = _to_internal_length(top)
        left_inches = to_inche(internal_left).value
        top_inches = to_inche(internal_top).value

        if component["type"] == "text":
            # レイアウト情報を一度だけ取得
            layout = component.get("layout", {}) or {}
            width = layout.get("width")
            height = layout.get("height")

            internal_width = _to_internal_length(width) if width else _Inch(3)
            internal_height = _to_internal_length(height) if height else _Inch(1)

            shape = slide_obj.shapes.add_textbox(
                pptx.util.Inches(left_inches),
                pptx.util.Inches(top_inches),
                pptx.util.Inches(to_inche(internal_width).value),
                pptx.util.Inches(to_inche(internal_height).value),
            )
            text_frame = shape.text_frame
            text_frame.text = component["text"]
            if size := component.get("size"):
                internal_size = _to_internal_length(size)
                text_frame.paragraphs[0].font.size = pptx.util.Pt(
                    to_point(internal_size).value
                )
            if bold := component.get("bold"):
                text_frame.paragraphs[0].font.bold = bold
            if italic := component.get("italic"):
                text_frame.paragraphs[0].font.italic = italic

        elif component["type"] == "image":
            # 画像のサイズ情報を一度だけ取得
            width = component.get("width")
            height = component.get("height")

            internal_width = _to_internal_length(width) if width else None
            internal_height = _to_internal_length(height) if height else None

            shape = slide_obj.shapes.add_picture(
                component["path"],
                pptx.util.Inches(left_inches),
                pptx.util.Inches(top_inches),
                pptx.util.Inches(to_inche(internal_width).value)
                if internal_width
                else None,
                pptx.util.Inches(to_inche(internal_height).value)
                if internal_height
                else None,
            )

        elif component["type"] == "chart":
            chart_data = ChartData()
            data = component["data"]
            if data and "category" in data[0]:
                chart_data.categories = [item["category"] for item in data]
            else:
                chart_data.categories = [f"Item {i + 1}" for i in range(len(data))]
            chart_data.add_series("Series 1", [item.get("value", 0) for item in data])

            # レイアウト情報を一度だけ取得
            layout = component.get("layout", {}) or {}
            width = layout.get("width")
            height = layout.get("height")

            internal_width = _to_internal_length(width) if width else _Inch(6)
            internal_height = _to_internal_length(height) if height else _Inch(4)

            x, y = pptx.util.Inches(left_inches), pptx.util.Inches(top_inches)
            cx = pptx.util.Inches(to_inche(internal_width).value)
            cy = pptx.util.Inches(to_inche(internal_height).value)

            shape = slide_obj.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED
                if component["chart_type"] == "bar"
                else XL_CHART_TYPE.LINE,
                x,
                y,
                cx,
                cy,
                chart_data,
            )

        elif component["type"] == "table":
            # レイアウト情報を一度だけ取得
            layout = component.get("layout", {})
            width = layout.get("width")
            height = layout.get("height")

            internal_width = _to_internal_length(width) if width else _Inch(6)
            internal_height = _to_internal_length(height) if height else _Inch(2)

            table = slide_obj.shapes.add_table(
                component["rows"],
                component["cols"],
                pptx.util.Inches(left_inches),
                pptx.util.Inches(top_inches),
                pptx.util.Inches(to_inche(internal_width).value),
                pptx.util.Inches(to_inche(internal_height).value),
            ).table

            for i, row in enumerate(component["data"]):
                for j, cell in enumerate(row):
                    table_cell = table.cell(i, j)
                    table_cell.text = cell.text

                    if size := cell.get("size"):
                        internal_size = _to_internal_length(size)
                        table_cell.text_frame.paragraphs[0].font.size = pptx.util.Pt(
                            to_point(internal_size).value
                        )
                    if bold := cell.get("bold"):
                        table_cell.text_frame.paragraphs[0].font.bold = bold
                    if italic := cell.get("italic"):
                        table_cell.text_frame.paragraphs[0].font.italic = italic
                    if color := cell.get("color"):
                        table_cell.text_frame.paragraphs[
                            0
                        ].font.color.rgb = RGBColor.from_string(color)
                    if background := cell.get("background"):
                        table_cell.fill.solid()
                        table_cell.fill.fore_color.rgb = RGBColor.from_string(
                            background
                        )

                    table_cell.text_frame.paragraphs[0].alignment = cell.get("align")

        if layout := component.get("layout"):
            self._apply_layout(shape, layout)

    def _add_container(
        self, slide_obj, container: Container, left: Length, top: Length
    ):
        """Add a container to a slide

        Args:
            slide_obj: Target slide
            container (Container): Container to add
            left (Length): Position from left edge
            top (Length): Position from top edge
        """
        current_left = left
        current_top = top

        for component in container["components"]:
            self._add_component(slide_obj, component, current_left, current_top)

            if container["layout"]["direction"] == "row":
                current_left = (current_left[0] + 2, current_left[1])  # Default spacing
            else:
                current_top = (current_top[0] + 1, current_top[1])  # Default spacing

    def build(self) -> "Presentation":
        """Build the presentation

        Returns:
            Presentation: Built presentation
        """
        for slide in self.slides:
            slide_layout = self.presentation.slide_layouts[slide["layout"].value]
            slide_obj = self.presentation.slides.add_slide(slide_layout)

            if title := slide.get("title"):
                title_shape = slide_obj.shapes.title
                title_shape.text = title["text"]
                if size := title.get("size"):
                    title_shape.text_frame.paragraphs[0].font.size = pptx.util.Pt(
                        to_point(size).value
                    )
                if bold := title.get("bold"):
                    title_shape.text_frame.paragraphs[0].font.bold = bold
                if italic := title.get("italic"):
                    title_shape.text_frame.paragraphs[0].font.italic = italic

            if containers := slide.get("containers"):
                for container in containers:
                    self._add_container(
                        slide_obj, container, (1, "in"), (2, "in")
                    )  # Default position

        return self.presentation


def image(
    path: str | pathlib.Path,
    width: Length | None = None,
    height: Length | None = None,
    layout: Layout | None = None,
) -> Image:
    """Create an image component with type field automatically set

    Args:
        path (str | pathlib.Path): Path to image file
        width (Length | None, optional): Width. Defaults to None.
        height (Length | None, optional): Height. Defaults to None.
        layout (Layout | None, optional): Layout settings. Defaults to None.

    Returns:
        Image: Created image component
    """
    return {
        "type": "image",
        "path": path,
        "width": width,
        "height": height,
        "layout": layout,
    }


def text(
    text: str,
    size: Length | None = None,
    bold: bool = False,
    italic: bool = False,
    color: str | None = None,
    layout: Layout | None = None,
) -> Text:
    """Create a text component with type field automatically set

    Args:
        text (str): Text content
        size (Length | None, optional): Font size. Defaults to None.
        bold (bool, optional): Whether text is bold. Defaults to False.
        italic (bool, optional): Whether text is italic. Defaults to False.
        color (str | None, optional): Text color. Defaults to None.
        layout (Layout | None, optional): Layout settings. Defaults to None.

    Returns:
        Text: Created text component
    """
    return {
        "type": "text",
        "text": text,
        "size": size,
        "bold": bold,
        "italic": italic,
        "color": color,
        "layout": layout,
    }


def chart(
    chart_type: str,
    data: list[dict[str, Any]],
    width: Length | None = None,
    height: Length | None = None,
    layout: Layout | None = None,
) -> Chart:
    """Create a chart component with type field automatically set

    Args:
        chart_type (str): Chart type ("bar", "line", "pie", etc.)
        data (list[dict[str, Any]]): Chart data
        width (Length | None, optional): Width. Defaults to None.
        height (Length | None, optional): Height. Defaults to None.
        layout (Layout | None, optional): Layout settings. Defaults to None.

    Returns:
        Chart: Created chart component
    """
    return {
        "type": "chart",
        "chart_type": chart_type,
        "data": data,
        "width": width,
        "height": height,
        "layout": layout,
    }


def table(
    rows: int,
    cols: int,
    data: list[list[TableCell]],
    width: Length | None = None,
    height: Length | None = None,
    layout: Layout | None = None,
) -> Table:
    """Create a table component with type field automatically set

    Args:
        rows (int): Number of rows
        cols (int): Number of columns
        data (list[list[TableCell]]): Table data
        width (Length | None, optional): Width. Defaults to None.
        height (Length | None, optional): Height. Defaults to None.
        layout (Layout | None, optional): Layout settings. Defaults to None.

    Returns:
        Table: Created table component
    """
    return {
        "type": "table",
        "rows": rows,
        "cols": cols,
        "data": data,
        "width": width,
        "height": height,
        "layout": layout,
    }


def layout(
    type: LayoutType = LayoutType.FLEX,
    direction: str = "row",
    align: Align = Align.START,
    justify: Justify = Justify.START,
    gap: Length = (0.1, "in"),
    padding: dict[str, Length] | None = None,
    width: Length | None = None,
    height: Length | None = None,
) -> Layout:
    """Create a layout with default values

    Args:
        type (LayoutType, optional): Layout type. Defaults to LayoutType.FLEX.
        direction (str, optional): Layout direction ("row" or "column"). Defaults to "row".
        align (Align, optional): Element alignment. Defaults to Align.START.
        justify (Justify, optional): Element justification. Defaults to Justify.START.
        gap (Length, optional): Gap between elements. Defaults to (0.1, "in").
        padding (dict[str, Length] | None, optional): Padding (top, right, bottom, left). Defaults to None.
        width (Length | None, optional): Width. Defaults to None.
        height (Length | None, optional): Height. Defaults to None.

    Returns:
        Layout: Created layout
    """
    return {
        "type": type,
        "direction": direction,
        "align": align,
        "justify": justify,
        "gap": gap,
        "padding": padding,
        "width": width,
        "height": height,
    }


def slide(
    layout: SlideLayout,
    title: Text | None = None,
    containers: list[Container] | None = None,
) -> Slide:
    """Create a slide with specified layout, title, and containers.

    Args:
        layout (SlideLayout): Layout type for the slide
        title (Text | None, optional): Title text component. Defaults to None.
        containers (list[Container] | None, optional): list of containers. Defaults to None.

    Returns:
        Slide: Created slide object
    """
    if containers is None:
        containers = []
    return {
        "layout": layout,
        "title": title,
        "containers": containers,
    }
