"""Table wrapper implementation."""

from typing import Literal, NotRequired, Self, TypeAlias, TypedDict

from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.shapes.graphfrm import GraphicFrame
from pptx.util import Pt

from pptxr.types._length import Length, LiteralLength

from .shape import Shape

DataFrame: TypeAlias = list[list[str]]


class TableCellStyle(TypedDict):
    """Table cell style properties."""

    text_align: NotRequired[Literal["left", "center", "right", "justify"]]
    vertical_align: NotRequired[Literal["top", "middle", "bottom"]]
    bold: NotRequired[bool]
    italic: NotRequired[bool]
    font_size: NotRequired[int]  # ポイントサイズで指定
    font_name: NotRequired[str]


class TableOptions(TypedDict):
    """Table properties."""

    left: Length | LiteralLength
    top: Length | LiteralLength
    width: Length | LiteralLength
    height: Length | LiteralLength
    cell_styles: NotRequired[list[list[TableCellStyle]]]
    first_row_header: NotRequired[bool]


class TableProps(TableOptions):
    """Table data."""

    type: Literal["table"]

    data: DataFrame


class Table(Shape[GraphicFrame]):
    """Table data class."""

    def __init__(
        self,
        pptx_obj: GraphicFrame,
        props: TableProps | None = None,
        /,
    ) -> None:
        self._pptx = pptx_obj

        if not props:
            return

        table = pptx_obj.table

        # Apply first row as header if specified
        if props.get("first_row_header"):
            table.first_row = True

        # Apply table data if provided
        if "data" in props:
            table_data = props["data"]
            for i, row in enumerate(table_data):
                if i < len(table.rows):
                    for j, cell_text in enumerate(row):
                        if j < len(table.columns):
                            table.cell(i, j).text = cell_text

        # Apply cell styles if provided
        if "cell_styles" in props:
            cell_styles = props["cell_styles"]
            for i, row_styles in enumerate(cell_styles):
                if i < len(table.rows):
                    for j, cell_style in enumerate(row_styles):
                        if j < len(table.columns):
                            cell = table.cell(i, j)

                            if "text_align" in cell_style:
                                align_map = {
                                    "left": PP_ALIGN.LEFT,
                                    "center": PP_ALIGN.CENTER,
                                    "right": PP_ALIGN.RIGHT,
                                    "justify": PP_ALIGN.JUSTIFY,
                                }
                                paragraph = cell.text_frame.paragraphs[0]
                                paragraph.alignment = align_map[
                                    cell_style["text_align"]
                                ]

                            if "vertical_align" in cell_style:
                                valign_map = {
                                    "top": MSO_VERTICAL_ANCHOR.TOP,
                                    "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
                                    "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
                                }
                                cell.text_frame.vertical_anchor = valign_map[
                                    cell_style["vertical_align"]
                                ]

                            # Apply text formatting
                            if any(
                                key in cell_style
                                for key in ["bold", "italic", "font_size", "font_name"]
                            ):
                                paragraph = cell.text_frame.paragraphs[0]
                                run = (
                                    paragraph.runs[0]
                                    if paragraph.runs
                                    else paragraph.add_run()
                                )

                                if "bold" in cell_style:
                                    run.font.bold = cell_style["bold"]

                                if "italic" in cell_style:
                                    run.font.italic = cell_style["italic"]

                                if "font_size" in cell_style:
                                    run.font.size = Pt(cell_style["font_size"])

                                if "font_name" in cell_style:
                                    run.font.name = cell_style["font_name"]

    def to_pptx(self) -> GraphicFrame:
        """Convert to pptx table frame."""
        return self._pptx

    @classmethod
    def from_pptx(cls, pptx_obj: GraphicFrame) -> Self:
        """Create from pptx table frame."""
        return cls(pptx_obj)
