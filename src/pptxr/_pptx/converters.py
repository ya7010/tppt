"""Converter implementations for pptx objects."""

import os
from typing import IO, Self, assert_never

import pptx
import pptx.util
from pptx import Presentation as PptxPresentation
from pptx import presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.shapes.base import BaseShape as PptxShape
from pptx.slide import Slide as PptxSlide

from pptxr.abstract.presentation import (
    AbstractPresentation,
    AbstractShape,
    AbstractSlide,
    PresentationFactory,
)
from pptxr.types import FilePath, Length, LiteralLength
from pptxr.types.length import Centimeter, Inch, Point, to_length

from .types import PptxConvertible
from .types import PptxPresentation as PptxPres

# pyright: ignore


class Shape(AbstractShape, PptxConvertible[PptxShape]):
    """Shape wrapper with type safety."""

    def __init__(self, pptx_shape: PptxShape) -> None:
        """Initialize shape."""
        self._pptx = pptx_shape

    def get_text(self) -> str:
        """Get shape text."""
        return self._pptx.text  # type: ignore

    def set_text(self, text: str) -> None:
        """Set shape text."""
        self._pptx.text = text  # type: ignore

    def to_pptx(self) -> PptxShape:
        """Convert to pptx shape."""
        return self._pptx

    @classmethod
    def from_pptx(cls, pptx_obj: PptxShape) -> Self:
        """Create from pptx shape."""
        return cls(pptx_obj)


class Slide(AbstractSlide, PptxConvertible[PptxSlide]):
    """Slide wrapper with type safety."""

    def __init__(self, pptx_slide: PptxSlide) -> None:
        """Initialize slide."""
        self._pptx = pptx_slide

    def get_shapes(self) -> list[AbstractShape]:
        """Get all shapes in the slide."""
        return [Shape(shape) for shape in self._pptx.shapes]

    def add_shape(
        self,
        shape_type: str,  # type: ignore
        left: Length | LiteralLength,
        top: Length | LiteralLength,
        width: Length | LiteralLength,
        height: Length | LiteralLength,
    ) -> AbstractShape:
        """Add a shape to the slide."""
        shape = self._pptx.shapes.add_shape(
            to_pptx_shape_type(shape_type),
            to_pptx_length(left),
            to_pptx_length(top),
            to_pptx_length(width),
            to_pptx_length(height),
        )
        return Shape(shape)

    def get_title(self) -> AbstractShape | None:
        """Get slide title shape."""
        if self._pptx.shapes.title:
            return Shape(self._pptx.shapes.title)
        return None

    def to_pptx(self) -> PptxSlide:
        """Convert to pptx slide."""
        return self._pptx

    @classmethod
    def from_pptx(cls, pptx_obj: PptxSlide) -> Self:
        """Create from pptx slide."""
        return cls(pptx_obj)


class Presentation(AbstractPresentation, PptxConvertible[presentation.Presentation]):
    """Presentation wrapper with type safety."""

    def __init__(
        self, pptx_presentation: presentation.Presentation | None = None
    ) -> None:
        """Initialize presentation."""
        if pptx_presentation is None:
            self._presentation = PptxPresentation()
        else:
            self._presentation = pptx_presentation

    def get_slides(self) -> list[AbstractSlide]:
        """Get all slides in the presentation."""
        return [Slide(slide) for slide in self._presentation.slides]

    def add_slide(self, layout_type: str) -> AbstractSlide:
        """Add a slide with specified layout."""
        layout_map = {
            "TITLE": 0,
            "TITLE_AND_CONTENT": 1,
            "SECTION_HEADER": 2,
            "TWO_CONTENT": 3,
            "COMPARISON": 4,
            "TITLE_ONLY": 5,
            "BLANK": 6,
            "CONTENT_WITH_CAPTION": 7,
            "PICTURE_WITH_CAPTION": 8,
            "TITLE_AND_VERTICAL_TEXT": 9,
            "VERTICAL_TITLE_AND_TEXT": 10,
        }
        layout = self._presentation.slide_layouts[layout_map[layout_type]]
        slide = self._presentation.slides.add_slide(layout)
        return Slide(slide)

    def save(self, file: FilePath | IO[bytes]) -> None:
        """Save presentation to file."""
        save_presentation(self._presentation, file)

    def to_pptx(self) -> presentation.Presentation:
        """Convert to pptx presentation."""
        return self._presentation

    @classmethod
    def from_pptx(cls, pptx_obj: presentation.Presentation) -> Self:
        """Create from pptx presentation."""
        return cls(pptx_obj)


class PptxPresentationFactory(PresentationFactory):
    """Factory for creating PPTX presentations."""

    def create_presentation(self) -> AbstractPresentation:
        """Create a new presentation."""
        return Presentation()

    def load_presentation(self, path: str) -> AbstractPresentation:
        """Load presentation from file."""
        return Presentation(PptxPresentation(path))


def to_pptx_length(length: Length | LiteralLength) -> pptx.util.Length:
    """Convert pptxr length to pptx length."""
    if isinstance(length, tuple):
        length = to_length(length)

    match length:
        case Inch():
            return pptx.util.Inches(length.value)
        case Centimeter():
            return pptx.util.Cm(length.value)
        case Point():
            return pptx.util.Pt(length.value)
        case _:
            assert_never(length)


def to_pptx_shape_type(shape_type: str) -> MSO_AUTO_SHAPE_TYPE:
    """Convert shape type string to MSO_AUTO_SHAPE_TYPE."""
    try:
        return getattr(MSO_AUTO_SHAPE_TYPE, shape_type)
    except AttributeError:
        raise ValueError(f"Invalid shape type: {shape_type}")


def save_presentation(presentation: PptxPres, file: FilePath | IO[bytes]) -> None:
    """Save presentation to file."""
    if isinstance(file, os.PathLike):
        file = os.fspath(file)
    presentation.save(file)


# Keep imports for type checking
__all__ = []
if False:
    PptxPresentation
    PptxShape
    PptxSlide
