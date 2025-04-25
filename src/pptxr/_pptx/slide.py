"""Slide wrapper implementation."""

from typing import TYPE_CHECKING, Self, Unpack, assert_never, cast

from pptx.slide import Slide as PptxSlide
from pptx.slide import SlideLayout as PptxSlideLayout

from pptxr.exception import SlideLayoutIndexError

from .converter import PptxConvertible, to_pptx_length
from .picture import Picture, PictureData, PictureProps
from .shape import Shape
from .table import Table, TableData, TableProps
from .text import Text, TextData, TextProps
from .title import Title

if TYPE_CHECKING:
    from .presentation import PresentationBuilder


class Slide(PptxConvertible[PptxSlide]):
    """Slide wrapper with type safety."""

    def __init__(self, pptx_slide: PptxSlide) -> None:
        """Initialize slide."""
        self._pptx: PptxSlide = pptx_slide

    @property
    def shapes(self) -> list[Shape]:
        """Get all shapes in the slide."""
        return [Shape(shape) for shape in self._pptx.shapes]

    @property
    def title(self) -> Title | None:
        """Get slide title shape."""
        if title := self._pptx.shapes.title:
            return Title(title)
        return None

    def to_pptx(self) -> PptxSlide:
        """Convert to pptx slide."""
        return cast(PptxSlide, self._pptx)

    @classmethod
    def from_pptx(cls, pptx_obj: PptxSlide) -> Self:
        """Create from pptx slide."""
        return cls(pptx_obj)


class SlideBuilder:
    """Slide builder."""

    def __init__(
        self,
        slide_layout: PptxSlideLayout | int = 0,
    ) -> None:
        self._slide_layout = slide_layout
        self._data: list[TextData | PictureData | TableData] = []

    def text(self, text: str, /, **kwargs: Unpack[TextProps]) -> Self:
        self._data.append(TextData(type="text", text=text, **kwargs))

        return self

    def picture(self, path: str, /, **kwargs: Unpack[PictureProps]) -> Self:
        self._data.append(PictureData(type="picture", **kwargs))

        return self

    def table(self, rows: int, cols: int, **kwargs: Unpack[TableProps]) -> Self:
        table_data: TableData = {"type": "table", "rows": rows, "cols": cols, **kwargs}  # type: ignore
        self._data.append(table_data)

        return self

    def _build(
        self,
        builder: "PresentationBuilder",
    ) -> Slide:
        if isinstance(self._slide_layout, int):
            try:
                slide_layout = builder._pptx.slide_layouts[self._slide_layout]
            except IndexError:
                raise SlideLayoutIndexError(
                    self._slide_layout,
                    builder._pptx.slide_layouts,
                )
        else:
            slide_layout = self._slide_layout

        slide = builder._pptx.slides.add_slide(slide_layout)

        for data in self._data:
            if data["type"] == "text":
                Text(
                    slide.shapes.add_textbox(
                        to_pptx_length(data["left"]),
                        to_pptx_length(data["top"]),
                        to_pptx_length(data["width"]),
                        to_pptx_length(data["height"]),
                    ),
                    data,
                )
            elif data["type"] == "picture":
                Picture(
                    slide.shapes.add_picture(
                        data["image_file"],
                        to_pptx_length(data["left"]),
                        to_pptx_length(data["top"]),
                        to_pptx_length(data.get("width")),
                        to_pptx_length(data.get("height")),
                    ),
                    data,
                )
            elif data["type"] == "table":
                Table(
                    slide.shapes.add_table(
                        data["rows"],
                        data["cols"],
                        to_pptx_length(data["left"]),
                        to_pptx_length(data["top"]),
                        to_pptx_length(data["width"]),
                        to_pptx_length(data["height"]),
                    ),
                    data,
                )
            else:
                assert_never(data)

        return Slide(slide)
