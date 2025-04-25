"""Presentation wrapper implementation."""

from pathlib import Path
from typing import IO, Optional, Protocol, TypeVar, Union, cast, runtime_checkable

from pptx import Presentation as PptxPresentation
from pptx.presentation import Presentation as PptxPresentationType
from pptx.shapes.base import BaseShape as PptxShape
from pptx.slide import Slide as PptxSlide

from pptxr._pptx.converters import save_presentation
from pptxr._pptx.slide import Slide
from pptxr.types import FilePath

T = TypeVar("T", bound="PptxConvertible")


@runtime_checkable
class PptxConvertible(Protocol):
    """Protocol for objects that can be converted to and from pptx objects."""

    def to_pptx(self) -> PptxPresentation | PptxSlide | PptxShape:
        """Convert to pptx object."""
        ...

    @classmethod
    def from_pptx(
        cls, pptx_obj: PptxPresentation | PptxSlide | PptxShape
    ) -> "PptxConvertible":
        """Create from pptx object."""
        ...


class Presentation(PptxConvertible):
    """Presentation wrapper with type safety."""

    def __init__(self, pptx_presentation: PptxPresentationType | None = None) -> None:
        """Initialize presentation."""
        if pptx_presentation is None:
            self._presentation = PptxPresentation()
        else:
            self._presentation = pptx_presentation

    def get_slides(self) -> list[Slide]:
        """Get all slides in the presentation."""
        return [Slide(slide) for slide in self._presentation.slides]

    def add_slide(self, layout_type: str) -> Slide:
        """Add a slide with specified layout."""
        layout_map = {
            "TITLE": 0,
            "TITLE_AND_CONTENT": 1,
            "SECTION_HEADER": 2,
            "TWO_CONTENT": 3,
            "COMPARISON": 4,
            "TITLE_ONLY": 5,
            "BLANK": 6,
            "CONTENT_WITH_CAPTION": 7,
            "PICTURE_WITH_CAPTION": 8,
            "TITLE_AND_VERTICAL_TEXT": 9,
            "VERTICAL_TITLE_AND_TEXT": 10,
        }
        layout = self._presentation.slide_layouts[layout_map[layout_type]]
        slide = self._presentation.slides.add_slide(layout)
        return Slide(slide)

    def save(self, file: FilePath | IO[bytes]) -> None:
        """Save presentation to file."""
        save_presentation(self._presentation, file)

    def to_pptx(self) -> PptxPresentation:
        """Convert to pptx presentation."""
        return cast(PptxPresentation, self._presentation)

    @classmethod
    def from_pptx(cls, pptx_obj: PptxPresentation) -> "Presentation":
        """Create from pptx presentation."""
        if not isinstance(pptx_obj, PptxPresentationType):
            raise TypeError(f"Expected PptxPresentation, got {type(pptx_obj)}")
        return cls(pptx_obj)


class PptxPresentationFactory:
    """A factory class for creating PowerPoint presentations."""

    def __init__(self, template_path: Optional[Union[str, Path]] = None) -> None:
        """Initialize a new presentation factory.

        Args:
            template_path: The path to the template file.
        """
        self._template_path = template_path
        self._presentation = (
            PptxPresentation(template_path) if template_path else PptxPresentation()
        )

    @classmethod
    def create(
        cls, template_path: Optional[Union[str, Path]] = None
    ) -> "PptxPresentationFactory":
        """Create a new presentation factory.

        Args:
            template_path: The path to the template file.

        Returns:
            A new presentation factory.
        """
        return cls(template_path)

    def add_slide(self, layout_type: str):
        """Add a slide to the presentation.

        Args:
            layout_type: The layout type of the slide.

        Returns:
            The added slide.
        """
        return self._presentation.slides.add_slide(self._presentation.slide_layouts[0])

    def save(self, path: Union[str, Path]) -> None:
        """Save the presentation to a file.

        Args:
            path: The path to save the presentation to.
        """
        self._presentation.save(str(path))
