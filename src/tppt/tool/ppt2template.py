"""Tool to generate slide master and layout type definitions from PowerPoint."""

import argparse
import os
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation as PptxPresentation

import tppt


@dataclass
class PlaceholderInfo:
    """Information about a placeholder."""

    name: str
    type: int
    field_name: Optional[str] = None
    field_type: Optional[str] = None
    required: bool = False
    idx: Optional[int] = None


@dataclass
class LayoutInfo:
    """Information about a slide layout."""

    name: str
    placeholders: List[PlaceholderInfo]
    class_name: Optional[str] = None


@dataclass
class MasterInfo:
    """Information about a slide master."""

    name: str
    layouts: List[LayoutInfo]
    class_name: Optional[str] = None


def extract_master_info(tree: Dict[str, Any]) -> MasterInfo:
    """Extract information about the slide master and its layouts."""

    master_name = Path(tree.get("name", "custom_template")).stem

    layouts_info = []

    # Extract layout information from slide masters
    if "slide_masters" in tree and tree["slide_masters"]:
        slide_master = tree["slide_masters"][0]
        if "slide_layouts" in slide_master:
            for layout in slide_master["slide_layouts"]:
                layout_name = layout.get("name", "")
                if layout_name:
                    placeholders = [
                        PlaceholderInfo(
                            name=p.get("name", ""),
                            type=p.get("placeholder_type", 0),
                            idx=i,
                        )
                        for i, p in enumerate(layout.get("placeholders", []))
                        if p.get("name")
                    ]
                    layouts_info.append(
                        LayoutInfo(name=layout_name, placeholders=placeholders)
                    )

    # Fallback if no layouts were extracted
    if not layouts_info:
        layout_map = {}
        for slide in tree.get("slides", []):
            layout_name = slide.get("slide_layout_name", "")
            if layout_name and layout_name not in layout_map:
                placeholders = [
                    PlaceholderInfo(
                        name=p.get("name", ""), type=p.get("placeholder_type", 0), idx=i
                    )
                    for i, p in enumerate(slide.get("placeholders", []))
                    if p.get("name")
                ]
                layout_map[layout_name] = LayoutInfo(
                    name=layout_name, placeholders=placeholders
                )
        layouts_info = list(layout_map.values())

    return MasterInfo(name=master_name, layouts=layouts_info)


def clean_field_name(name: str) -> str:
    """Clean a field name to make it suitable for a Python identifier.

    Removes numeric suffixes, converts to snake_case, and handles special characters.
    """
    # Remove numeric suffixes (e.g., "Title 1" -> "Title")
    name = re.sub(r"\s+\d+$", "", name)

    # Convert to lowercase and replace spaces/special chars with underscores
    name = name.lower().replace(" ", "_").replace("-", "_").replace(".", "_")

    # Remove "placeholder" suffix if present
    name = name.replace("_placeholder", "")

    # Ensure the name is a valid Python identifier
    if not name or not name[0].isalpha() and name[0] != "_":
        name = f"placeholder_{name}"

    # Replace any remaining invalid characters
    name = re.sub(r"[^a-z0-9_]", "_", name)

    return name


def analyze_layout(layout: LayoutInfo) -> LayoutInfo:
    """Analyze layout and determine class name and field information."""

    # Generate class name in PascalCase
    layout_name_parts = layout.name.replace("-", " ").split()
    class_name = (
        "Custom" + "".join(part.capitalize() for part in layout_name_parts) + "Layout"
    )

    # Process placeholders to ensure no duplicates and proper field names
    processed_placeholders = []
    seen_field_names = set()
    placeholder_type_map = {}
    layout_name_lower = layout.name.lower()

    # 同じ種類のプレースホルダー名をカウント (例：contentという名前が何個あるか)
    placeholder_base_name_counts = {}

    # First pass - identify placeholder types and count base names
    for ph in layout.placeholders:
        if ph.type not in placeholder_type_map:
            placeholder_type_map[ph.type] = []
        placeholder_type_map[ph.type].append(ph)

        # クリーニングされた基本名をカウント
        base_name = clean_field_name(ph.name)
        if ph.type in [2, 7]:  # content/bodyタイプのプレースホルダー
            if "content" in layout_name_lower or "text" in layout_name_lower:
                base_name = "content"
            else:
                base_name = "body"

        placeholder_base_name_counts[base_name] = (
            placeholder_base_name_counts.get(base_name, 0) + 1
        )

    # 同じ種類のプレースホルダーに対してベース名を決める
    type_base_names = {}
    for ph_type, phs in placeholder_type_map.items():
        # タイプごとにベース名を決定
        if len(phs) == 0:
            continue

        sample_ph = phs[0]
        base_name = ""

        # タイプに基づいてベース名を決定
        if ph_type == 1:  # Title
            base_name = "title"
        elif ph_type == 2:  # Body/Content
            if "content" in layout_name_lower or "text" in layout_name_lower:
                base_name = "content"
            else:
                base_name = "body"
        elif ph_type == 3:  # CenteredTitle
            base_name = "title"
        elif ph_type == 4:  # Subtitle
            base_name = "subtitle"
        elif ph_type == 7:  # Chart
            if "chart" in sample_ph.name.lower():
                base_name = "chart"
            else:
                base_name = "content"
        elif ph_type == 8:  # Table
            base_name = "table"
        elif ph_type == 13:  # SlideNumber
            base_name = "slide_number"
        elif ph_type == 15:  # Footer
            base_name = "footer"
        elif ph_type == 16:  # Date
            base_name = "date"
        elif ph_type == 18:  # Picture
            base_name = "picture"
        elif ph_type == 19:  # VerticalTitle
            base_name = "vertical_title"
        elif ph_type == 20:  # VerticalBody
            base_name = "vertical_text"
        else:
            # その他のタイプはプレースホルダー名をクリーニングして使用
            base_name = clean_field_name(sample_ph.name)

        # すべてのタイプのプレースホルダーについて処理
        if ph_type not in [
            1,
            3,
            4,
            13,
            15,
            16,
            19,
            20,
        ]:  # タイトル、サブタイトル、日付、フッターなどは連番を付けない
            # 特定のレイアウトタイプの特別処理
            if (
                "two content" in layout_name_lower
                and ph_type in [2, 7]
                and len(phs) == 2
            ):
                type_base_names[ph_type] = ["left_content", "right_content"]
            elif (
                "comparison" in layout_name_lower
                and ph_type in [2, 7]
                and len(phs) >= 4
            ):
                special_names = [
                    "left_title" if "title" in phs[0].name.lower() else "left_content",
                    "left_body" if "body" in phs[1].name.lower() else "left_content",
                    "right_title"
                    if "title" in phs[2].name.lower()
                    else "right_content",
                    "right_body" if "body" in phs[3].name.lower() else "right_content",
                ]
                type_base_names[ph_type] = special_names

                # 残りがあれば、content5, content6...
                for i in range(4, len(phs)):
                    type_base_names[ph_type].append(f"content{i + 1}")
            else:
                # 同じ種類のプレースホルダーが複数ある場合のみ連番をつける
                if len(phs) > 1:
                    base_name_without_numbers = re.sub(
                        r"\d+$", "", base_name
                    )  # 名前から末尾の数字を取り除く
                    type_base_names[ph_type] = [
                        f"{base_name_without_numbers}{i + 1}" for i in range(len(phs))
                    ]
                else:
                    # 単一のプレースホルダーの場合は連番なし
                    type_base_names[ph_type] = [base_name]
        else:
            # 日付、フッター、スライド番号などには連番を付けない
            type_base_names[ph_type] = [base_name]

    # Second pass - 実際に各プレースホルダーにフィールド名を割り当てる
    processed_placeholders = []

    for ph in layout.placeholders:
        if not ph.name:
            continue

        placeholder_type = ph.type

        # このタイプのプレースホルダーの何番目か
        same_type_phs = placeholder_type_map.get(placeholder_type, [])
        idx = same_type_phs.index(ph)

        # フィールド名の決定
        if placeholder_type in type_base_names and idx < len(
            type_base_names[placeholder_type]
        ):
            field_name = type_base_names[placeholder_type][idx]
        else:
            # フォールバック: 元の名前をクリーニング
            field_name = clean_field_name(ph.name)

        # 名前の重複を避ける
        if field_name in seen_field_names:
            # Content with Captionレイアウトの場合は特別処理
            if layout.name == "Content with Caption" and placeholder_type == 2:
                # ここでは何もしない（上ですでに正しい名前を設定）
                pass
            else:
                counter = 1
                while f"{field_name}{counter}" in seen_field_names:
                    counter += 1
                field_name = f"{field_name}{counter}"

        seen_field_names.add(field_name)

        # Determine field type
        required = False

        if "date" in field_name or placeholder_type == 16:  # Date
            field_type = "tppt.Placeholder[datetime.date | None]"
        elif "number" in field_name or placeholder_type == 13:  # Slide number
            field_type = 'tppt.Placeholder[Literal["‹#›"] | int | None]'
        elif "title" in field_name or placeholder_type in [
            1,
            3,
            19,
        ]:  # Title, CenteredTitle, VerticalTitle
            field_type = "tppt.Placeholder[str]"
            required = True
        elif "subtitle" in field_name or placeholder_type == 4:  # Subtitle
            field_type = "tppt.Placeholder[str | None]"
        elif (
            "picture" in field_name or "image" in field_name or placeholder_type == 18
        ):  # Picture
            field_type = "tppt.Placeholder[tppt.types.FilePath | None]"
        elif "chart" in field_name and placeholder_type == 7:  # Actual chart
            field_type = "tppt.Placeholder[tppt.types.FilePath | None]"
        elif "table" in field_name and placeholder_type == 8:  # Table
            field_type = "tppt.Placeholder[str | None]"
        elif placeholder_type in [
            2,
            7,
            20,
        ]:  # Body, Chart (used as content), VerticalBody
            field_type = "tppt.Placeholder[str | None]"
        elif "footer" in field_name or placeholder_type == 15:  # Footer
            field_type = "tppt.Placeholder[str | None]"
        else:
            field_type = "tppt.Placeholder[str | None]"

        new_ph = PlaceholderInfo(
            name=ph.name,
            type=ph.type,
            field_name=field_name,
            field_type=field_type,
            required=required,
            idx=ph.idx,
        )
        processed_placeholders.append(new_ph)

    return LayoutInfo(
        name=layout.name, placeholders=processed_placeholders, class_name=class_name
    )


def generate_layout_class(layout: LayoutInfo) -> str:
    """Generate Python code for a slide layout class."""

    class_name = layout.class_name or f"Custom{layout.name.replace(' ', '')}Layout"
    class_docstring = f'"""{layout.name} layout."""'

    # 重複するフィールド名を削除する
    unique_fields = {}
    for ph in layout.placeholders:
        if not ph.field_name or not ph.field_type:
            continue

        # 同じフィールド名が存在する場合は上書き（最後のものを使用）
        unique_fields[ph.field_name] = ph

    # 元のインデックス順を維持する（ソートしない）
    placeholder_definitions = []
    for ph in layout.placeholders:
        if not ph.field_name or not ph.field_type:
            continue

        # 重複を避けるため、最後に処理されたフィールドだけを出力
        if unique_fields.get(ph.field_name) is not ph:
            continue

        if ph.required:
            placeholder_definitions.append(f"    {ph.field_name}: {ph.field_type}")
        else:
            placeholder_definitions.append(
                f"    {ph.field_name}: {ph.field_type} = None"
            )

    if not placeholder_definitions:
        placeholder_definitions = ["    # No placeholders found"]

    return (
        f"class {class_name}(tppt.SlideLayout):\n"
        f"    {class_docstring}\n\n" + "\n\n".join(placeholder_definitions)
    )


def generate_master_class(master: MasterInfo) -> str:
    """Generate Python code for a slide master class."""

    # Convert master_name to PascalCase
    class_prefix = "".join(
        word.capitalize() for word in master.name.replace("-", "_").split("_")
    )
    class_name = f"{class_prefix}SlideMaster"

    layout_definitions = []
    for layout in master.layouts:
        layout_class_name = (
            layout.class_name or f"Custom{layout.name.replace(' ', '')}Layout"
        )
        layout_field_name = f"{layout.name.replace(' ', '')}Layout"
        layout_definitions.append(
            f"    {layout_field_name}: tppt.Layout[{layout_class_name}]"
        )

    if not layout_definitions:
        layout_definitions = ["    # No layouts found"]

    class_docstring = f'"""{master.name} slide master."""'

    return (
        f'@tppt.slide_master("{master.name}.pptx")\n'
        f"class {class_name}(tppt.SlideMaster):\n"
        f"    {class_docstring}\n\n" + "\n\n".join(layout_definitions)
    )


def analyze_slide_structure(pptx_path: tppt.types.FilePath) -> dict[str, Any]:
    """Analyzes slide master and layout information from a PowerPoint file.

    Args:
        pptx_path: Path to the PowerPoint file
    """
    pptx_path = Path(pptx_path)

    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    # Load the presentation
    presentation = tppt.Presentation.from_pptx(PptxPresentation(os.fspath(pptx_path)))

    # Get the tree structure
    tree = presentation.tree
    return tree


def generate_template_file(
    pptx_path: tppt.types.FilePath,
    *,
    output_path: tppt.types.FilePath | None = None,
) -> None:
    """Generates slide master and layout definitions from a PowerPoint file and saves them as a Python file.

    Args:
        pptx_path: Path to the PowerPoint file
        output_path: Path to save the generated Python file (default: prints to standard output)
    """
    pptx_path = Path(pptx_path)

    # Analyze slide structure
    tree = analyze_slide_structure(pptx_path)
    tree["name"] = str(pptx_path)

    # Extract and analyze information
    master_info = extract_master_info(tree)

    # Analyze each layout
    analyzed_layouts = []
    for layout in master_info.layouts:
        analyzed_layout = analyze_layout(layout)
        analyzed_layouts.append(analyzed_layout)

    master_info.layouts = analyzed_layouts

    # Generate code for imports
    imports = [
        "import datetime",
        "from typing import Literal",
        "",
        "import tppt",
        "",
    ]

    # Generate code for layout classes
    layout_classes = []
    for layout in master_info.layouts:
        layout_class = generate_layout_class(layout)
        layout_classes.append(layout_class)

    # Generate master class
    master_class = generate_master_class(master_info)

    # Combine everything
    content = (
        "\n".join(imports)
        + "\n\n"
        + "\n\n\n".join(layout_classes)
        + "\n\n\n"
        + master_class
    )

    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(content)
    else:
        print(content)


if __name__ == "__main__":
    try:
        from rich_argparse import RichHelpFormatter

        formatter_class = RichHelpFormatter
    except ImportError:
        formatter_class = argparse.HelpFormatter

    parser = argparse.ArgumentParser(
        description="Generate slide master and layout type definitions from PowerPoint",
        formatter_class=formatter_class,
    )
    parser.add_argument("pptx_path", help="Path to the PPTX file")
    parser.add_argument(
        "-o", "--output", default=None, help="Path to save the generated Python file"
    )

    args = parser.parse_args()

    generate_template_file(args.pptx_path, output_path=args.output)
