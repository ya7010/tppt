"""Presentation wrapper implementation."""

import os
from typing import IO, TYPE_CHECKING, Any, Callable, Generic, Self, overload

from pptx.presentation import Presentation as PptxPresentation

from tppt._pptx.tree import ppt2dict
from tppt.slide_layout import TpptSlideLayout
from tppt.slide_master import (
    DefaultSlideMaster,
    GenericTpptSlideMaster,
)
from tppt.types import FilePath

from .converter import PptxConvertible
from .slide import SlideBuilder
from .slide_master import SlideMaster

if TYPE_CHECKING:
    from tppt._pptx.slide import Slide


class Presentation(PptxConvertible[PptxPresentation]):
    """Presentation wrapper with type safety."""

    def __init__(
        self,
        presentation: PptxPresentation,
    ) -> None:
        """Initialize presentation."""
        self._pptx = presentation

    @property
    def slides(self) -> "list[Slide]":
        """Get the slides."""
        from tppt._pptx.slide import Slide

        return [Slide.from_pptx(slide) for slide in self._pptx.slides]

    @property
    def slide_master(self) -> SlideMaster:
        """
        Get the slide master.

        This tool supports only one slide master.
        """
        return SlideMaster.from_pptx(self._pptx.slide_masters[0])

    @property
    def tree(self) -> dict[str, Any]:
        """Get the node tree of the presentation."""
        return ppt2dict(self._pptx)

    @overload
    @classmethod
    def builder(
        cls,
    ) -> "PresentationBuilder[DefaultSlideMaster]": ...

    @overload
    @classmethod
    def builder(
        cls,
        slide_master: "type[GenericTpptSlideMaster]",
    ) -> "PresentationBuilder[GenericTpptSlideMaster]": ...

    @classmethod
    def builder(
        cls,
        slide_master=None,
    ):
        """Get a builder for the presentation."""

        if slide_master is None:
            slide_master = DefaultSlideMaster
        return PresentationBuilder(slide_master)

    def save(self, file: FilePath | IO[bytes]) -> None:
        """Save presentation to file."""
        if isinstance(file, os.PathLike):
            file = os.fspath(file)
        self._pptx.save(file)

    def to_pptx(self) -> PptxPresentation:
        """Convert to pptx presentation."""
        return self._pptx

    @classmethod
    def from_pptx(cls, pptx_obj: PptxPresentation) -> Self:
        """Create from pptx presentation."""
        return cls(pptx_obj)


class PresentationBuilder(Generic[GenericTpptSlideMaster]):
    """Builder for presentations."""

    def __init__(
        self,
        slide_master: "type[GenericTpptSlideMaster]",
    ) -> None:
        """Initialize the builder."""
        import pptx

        self._pptx = pptx.Presentation()
        self._slide_master = slide_master

    def slide(
        self,
        slide: SlideBuilder
        | Callable[[type[GenericTpptSlideMaster]], SlideBuilder | TpptSlideLayout],
        /,
    ) -> Self:
        """Add a slide to the presentation."""
        if isinstance(slide, SlideBuilder):
            slide._build(self)
        else:
            slide_layout = slide(self._slide_master)
            slide_builder = (
                slide_layout.builder()
                if isinstance(slide_layout, TpptSlideLayout)
                else slide_layout
            )
            slide_builder._build(self)

        return self

    def build(self) -> Presentation:
        """Build the presentation."""

        return Presentation(self._pptx)

    def save(self, file: FilePath | IO[bytes]) -> None:
        """Save the presentation to a file."""
        self.build().save(file)
