"""Slide wrapper implementation."""

import os
from typing import IO, TYPE_CHECKING, Any, Callable, Self, Unpack, cast

from pptx.slide import Slide as PptxSlide
from pptx.slide import SlideLayout as PptxSlideLayout

from tppt._pptx.placeholder import SlidePlaceholder
from tppt.exception import SlideLayoutIndexError
from tppt.types import FilePath

from .converter import PptxConvertible, to_pptx_length
from .shape import Shape
from .shape.picture import Picture, PictureData, PictureProps
from .shape.table import DataFrame, Table, TableData, TableProps, dataframe2list
from .shape.text import Text, TextData, TextProps
from .shape.title import Title

if TYPE_CHECKING:
    from .presentation import PresentationBuilder


class Slide(PptxConvertible[PptxSlide]):
    """Slide wrapper with type safety."""

    def __init__(self, pptx_slide: PptxSlide) -> None:
        """Initialize slide."""
        self._pptx: PptxSlide = pptx_slide

    @property
    def shapes(self) -> list[Shape]:
        """Get all shapes in the slide."""
        return [Shape(shape) for shape in self._pptx.shapes]

    @property
    def title(self) -> Title | None:
        """Get slide title shape."""

        if title := self._pptx.shapes.title:
            return Title(title)
        return None

    @property
    def placeholders(self) -> list[SlidePlaceholder]:
        """Get all placeholders in the slide."""
        return [
            SlidePlaceholder(
                placeholder,  # type: ignore
            )
            for placeholder in self._pptx.placeholders
        ]

    def to_pptx(self) -> PptxSlide:
        """Convert to pptx slide."""
        return cast(PptxSlide, self._pptx)

    @classmethod
    def from_pptx(cls, pptx_obj: PptxSlide) -> Self:
        """Create from pptx slide."""
        return cls(pptx_obj)


class SlideBuilder:
    """Slide builder."""

    def __init__(
        self,
        slide_layout: PptxSlideLayout | int = 0,
    ) -> None:
        self._slide_layout = slide_layout
        self._shape_registry: list[Callable[[PptxSlide], Shape[Any]]] = []

    def text(self, text: str, **kwargs: Unpack[TextProps]) -> Self:
        data = TextData(type="text", text=text, **kwargs)

        self._shape_registry.append(
            lambda slide: Text(
                slide.shapes.add_textbox(
                    to_pptx_length(data["left"]),
                    to_pptx_length(data["top"]),
                    to_pptx_length(data["width"]),
                    to_pptx_length(data["height"]),
                ),
                data,
            )
        )

        return self

    def picture(
        self, image_file: FilePath | IO[bytes], **kwargs: Unpack[PictureProps]
    ) -> Self:
        data = PictureData(type="picture", image_file=image_file, **kwargs)

        if isinstance(data["image_file"], os.PathLike):
            image_file = os.fspath(data["image_file"])
        else:
            image_file = data["image_file"]

        self._shape_registry.append(
            lambda slide: Picture(
                slide.shapes.add_picture(
                    image_file,
                    to_pptx_length(data["left"]),
                    to_pptx_length(data["top"]),
                    to_pptx_length(data.get("width")),
                    to_pptx_length(data.get("height")),
                ),
                data,
            )
        )

        return self

    def table(self, data: DataFrame, **kwargs: Unpack[TableProps]) -> Self:
        data = dataframe2list(data)
        rows, cols = len(data), len(data[0])
        table_data: TableData = {"type": "table", "data": data, **kwargs}

        self._shape_registry.append(
            lambda slide: Table(
                slide.shapes.add_table(
                    rows,
                    cols,
                    to_pptx_length(table_data["left"]),
                    to_pptx_length(table_data["top"]),
                    to_pptx_length(table_data["width"]),
                    to_pptx_length(table_data["height"]),
                ),
                table_data,
            )
        )
        return self

    def _build(
        self,
        builder: "PresentationBuilder",
    ) -> Slide:
        if isinstance(self._slide_layout, int):
            try:
                slide_layout = builder._pptx.slide_layouts[self._slide_layout]
            except IndexError:
                raise SlideLayoutIndexError(
                    self._slide_layout,
                    builder._pptx.slide_layouts,
                )
        else:
            slide_layout = self._slide_layout

        slide = builder._pptx.slides.add_slide(slide_layout)

        for register in self._shape_registry:
            register(slide)

        return Slide(slide)
