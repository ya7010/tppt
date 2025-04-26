import argparse
import json
import os
from typing import Any, Dict

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation as PptxPresentation


def color_format_to_dict(color_format: Any) -> Dict[str, Any]:
    """Convert color format to dictionary"""
    data = {}

    try:
        if hasattr(color_format, "rgb"):
            if color_format.rgb is not None:
                data["rgb"] = str(color_format.rgb)

        if hasattr(color_format, "theme_color"):
            if color_format.theme_color is not None:
                data["theme_color"] = str(color_format.theme_color)

        if hasattr(color_format, "brightness"):
            data["brightness"] = color_format.brightness
    except Exception:
        pass

    return data


def text_frame_to_dict(text_frame: Any) -> Dict[str, Any]:
    """Convert text frame to dictionary"""
    if not hasattr(text_frame, "paragraphs"):
        return {"text": str(text_frame) if text_frame else ""}

    result = {"text": text_frame.text, "paragraphs": []}

    for p in text_frame.paragraphs:
        paragraph = {"text": p.text, "level": p.level, "runs": []}

        for run in p.runs:
            run_data = {"text": run.text}

            if hasattr(run, "font"):
                font_data = {}
                if hasattr(run.font, "name") and run.font.name is not None:
                    font_data["name"] = run.font.name
                if hasattr(run.font, "size") and run.font.size is not None:
                    font_data["size"] = (
                        run.font.size.pt
                        if hasattr(run.font.size, "pt")
                        else run.font.size
                    )
                if hasattr(run.font, "bold") and run.font.bold is not None:
                    font_data["bold"] = run.font.bold
                if hasattr(run.font, "italic") and run.font.italic is not None:
                    font_data["italic"] = run.font.italic
                if hasattr(run.font, "underline") and run.font.underline is not None:
                    font_data["underline"] = run.font.underline

                if hasattr(run.font, "color") and run.font.color is not None:
                    font_data["color"] = color_format_to_dict(run.font.color)

                run_data["font"] = font_data

            paragraph["runs"].append(run_data)

        result["paragraphs"].append(paragraph)

    return result


def shape_to_dict(shape: Any) -> Dict[str, Any]:
    """Convert Shape object to dictionary"""
    shape_data = {
        "name": shape.name,
        "shape_id": shape.shape_id,
        "shape_type": str(shape.shape_type),
        "has_text_frame": shape.has_text_frame,
        "has_table": shape.has_table,
        "has_chart": shape.has_chart,
        "width": shape.width.pt if hasattr(shape.width, "pt") else shape.width,
        "height": shape.height.pt if hasattr(shape.height, "pt") else shape.height,
        "rotation": shape.rotation,
        "left": shape.left.pt if hasattr(shape.left, "pt") else shape.left,
        "top": shape.top.pt if hasattr(shape.top, "pt") else shape.top,
    }

    # プレースホルダーの場合
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        try:
            shape_data["placeholder_type"] = (
                shape.placeholder_format.type
                if hasattr(shape.placeholder_format, "type")
                else None
            )
            shape_data["placeholder_idx"] = (
                shape.placeholder_format.idx
                if hasattr(shape.placeholder_format, "idx")
                else None
            )
        except Exception:
            pass

    # テキストフレームがある場合
    if shape.has_text_frame:
        shape_data["text_frame"] = text_frame_to_dict(shape.text_frame)

    # テーブルがある場合
    if shape.has_table:
        table_data = {
            "rows": shape.table.rows.count,
            "columns": shape.table.columns.count,
            "cells": [],
        }

        for row_idx in range(shape.table.rows.count):
            for col_idx in range(shape.table.columns.count):
                cell = shape.table.cell(row_idx, col_idx)
                cell_data = {
                    "row": row_idx,
                    "column": col_idx,
                    "text_frame": text_frame_to_dict(cell.text_frame),
                }
                table_data["cells"].append(cell_data)

        shape_data["table"] = table_data

    # グループシェイプの場合
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        shape_data["shapes"] = [shape_to_dict(subshape) for subshape in shape.shapes]

    return shape_data


def placeholder_to_dict(placeholder: Any) -> Dict[str, Any]:
    """Convert placeholder to dictionary"""
    placeholder_data = shape_to_dict(placeholder)

    # プレースホルダー特有の情報を追加
    try:
        placeholder_data["placeholder_type"] = (
            placeholder.placeholder_format.type
            if hasattr(placeholder.placeholder_format, "type")
            else None
        )
        placeholder_data["placeholder_idx"] = (
            placeholder.placeholder_format.idx
            if hasattr(placeholder.placeholder_format, "idx")
            else None
        )
    except Exception:
        pass

    return placeholder_data


def slide_layout_to_dict(slide_layout: Any) -> Dict[str, Any]:
    """Convert slide layout to dictionary"""
    layout_data = {
        "name": slide_layout.name,
        "slide_master_id": id(slide_layout.slide_master)
        if hasattr(slide_layout, "slide_master")
        else None,
        "shapes": [shape_to_dict(shape) for shape in slide_layout.shapes],
        "placeholders": [
            placeholder_to_dict(placeholder)
            for placeholder in slide_layout.placeholders
        ],
    }

    return layout_data


def slide_master_to_dict(slide_master: Any) -> Dict[str, Any]:
    """Convert slide master to dictionary"""
    master_data = {
        "id": id(slide_master),
        "shapes": [shape_to_dict(shape) for shape in slide_master.shapes],
        "placeholders": [
            placeholder_to_dict(placeholder)
            for placeholder in slide_master.placeholders
        ],
        "slide_layouts": [],
    }

    # スライドレイアウトの情報を追加
    for layout in slide_master.slide_layouts:
        layout_dict = slide_layout_to_dict(layout)
        if slide_layouts := master_data.get("slide_layouts"):
            slide_layouts.append(layout_dict)  # type: ignore

    return master_data


def slide_to_dict(slide: Any) -> Dict[str, Any]:
    """Convert slide to dictionary"""
    slide_data = {
        "slide_id": slide.slide_id,
        "slide_layout_name": slide.slide_layout.name
        if hasattr(slide.slide_layout, "name")
        else None,
        "shapes": [shape_to_dict(shape) for shape in slide.shapes],
        "placeholders": [],
    }

    # プレースホルダーの情報を追加
    if hasattr(slide, "placeholders"):
        slide_data["placeholders"] = [
            placeholder_to_dict(placeholder) for placeholder in slide.placeholders
        ]

    # ノートの情報を追加
    if hasattr(slide, "notes_slide") and slide.notes_slide is not None:
        notes_placeholders = []
        if hasattr(slide.notes_slide, "placeholders"):
            notes_placeholders = [
                placeholder_to_dict(placeholder)
                for placeholder in slide.notes_slide.placeholders
            ]

        slide_data["notes_slide"] = {
            "shapes": [shape_to_dict(shape) for shape in slide.notes_slide.shapes],
            "placeholders": notes_placeholders,
        }

    return slide_data


def presentation_to_dict(ppt: PptxPresentation) -> Dict[str, Any]:
    """Convert presentation information to dictionary"""

    # スライドのサイズを安全に取得
    slide_width = None
    slide_height = None

    if slide_width := getattr(ppt, "slide_width"):
        if pt := getattr(slide_width, "pt"):
            slide_width = pt
        else:
            slide_width = slide_width

    if slide_height := getattr(ppt, "slide_height"):
        if pt := getattr(slide_height, "pt"):
            slide_height = pt
        else:
            slide_height = ppt.slide_height

    prs_data = {
        "slides_count": len(ppt.slides),
        "slide_masters_count": len(ppt.slide_masters),
        "slide_layouts_count": len(ppt.slide_layouts),
        "slide_width": slide_width,
        "slide_height": slide_height,
        "slides": [slide_to_dict(slide) for slide in ppt.slides],
        "slide_masters": [slide_master_to_dict(master) for master in ppt.slide_masters],
    }

    # ノートマスターの情報を追加
    if notes_master := getattr(ppt, "notes_master"):
        notes_placeholders = []
        if placeholders := getattr(notes_master, "placeholders"):
            notes_placeholders = [
                placeholder_to_dict(placeholder) for placeholder in placeholders
            ]

        prs_data["notes_master"] = {
            "shapes": [shape_to_dict(shape) for shape in ppt.notes_master.shapes],
            "placeholders": notes_placeholders,
        }

    return prs_data


def main():
    parser = argparse.ArgumentParser(
        description="Tool to convert PPTX file structure to JSON"
    )
    parser.add_argument("pptx_file", help="Path to the PPTX file to convert")
    parser.add_argument(
        "-o", "--output", help="Path to the output JSON file", default=None
    )
    parser.add_argument(
        "--indent", type=int, help="JSON indentation for formatting", default=2
    )
    args = parser.parse_args()

    if not os.path.exists(args.pptx_file):
        print(f"Error: File '{args.pptx_file}' not found.")
        return

    prs_data = presentation_to_dict(args.pptx_file)

    # JSONファイルに出力
    if args.output:
        output_path = args.output
    else:
        base_name = os.path.splitext(os.path.basename(args.pptx_file))[0]
        output_path = f"{base_name}_structure.json"

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(prs_data, f, ensure_ascii=False, indent=args.indent)

    print(f"Structure of '{args.pptx_file}' has been output to '{output_path}'.")
