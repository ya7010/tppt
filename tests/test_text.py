"""Tests for text module."""

from typing import cast

from pptx.shapes.autoshape import Shape as PptxShape
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN

import tppt


def test_text_with_options(output) -> None:
    """Test creating text with various formatting options."""
    # テキストとその書式設定オプション
    # プレゼンテーションの作成
    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .text(
                "サンプルテキスト1",
                left=(100, "pt"),
                top=(300, "pt"),
                width=(400, "pt"),
                height=(200, "pt"),
                size=(24, "pt"),
                bold=True,
                italic=True,
                color="#0C0",
                margin_bottom=(10, "pt"),
                margin_left=(10, "pt"),
                vertical_anchor=MSO_ANCHOR.MIDDLE,
                word_wrap=True,
                auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                alignment=PP_ALIGN.CENTER,
                level=1,
            )
            .text(
                "サンプルテキスト2",
                left=(100, "pt"),
                top=(100, "pt"),
                width=(400, "pt"),
                height=(200, "pt"),
                size=(24, "pt"),
                bold=True,
                italic=True,
                color="#0000BB99",
                margin_bottom=(10, "pt"),
                margin_left=(10, "pt"),
                vertical_anchor=MSO_ANCHOR.MIDDLE,
                word_wrap=True,
                auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                alignment=PP_ALIGN.CENTER,
                level=1,
            )
        )
        .build()
    )

    # プレゼンテーションを保存
    pptx_path = output / "text_with_options.pptx"
    presentation.save(pptx_path)


def test_text_default_options(output) -> None:
    """Test creating text with default options."""
    text_content = "デフォルト設定のテキスト"

    # デフォルト設定でテキストを作成
    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.TitleLayout(
                title="Title",
            )
            .builder()
            .text(
                text_content,
                left=(100, "pt"),
                top=(100, "pt"),
                width=(300, "pt"),
                height=(100, "pt"),
            )
        )
        .build()
    )

    # プレゼンテーションを保存
    pptx_path = output / "text_default_options.pptx"
    presentation.save(pptx_path)


def test_font_language_id(output) -> None:
    """Test Font.language_id get/set and chaining."""
    from pptx.enum.lang import MSO_LANGUAGE_ID

    from tppt.pptx.text.font import Font

    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .text(
                "Language test",
                left=(100, "pt"),
                top=(100, "pt"),
                width=(300, "pt"),
                height=(100, "pt"),
            )
        )
        .build()
    )

    # Access the font from the presentation
    pptx_slide = presentation.to_pptx().slides[0]
    shape = cast(PptxShape, pptx_slide.shapes[0])
    pptx_font = shape.text_frame.paragraphs[0].runs[0].font

    font = Font(pptx_font)

    # Test setter
    font.language_id = MSO_LANGUAGE_ID.JAPANESE
    assert font.language_id == MSO_LANGUAGE_ID.JAPANESE

    # Test chaining
    result = font.set_language_id(MSO_LANGUAGE_ID.ENGLISH_US)
    assert result is font
    assert font.language_id == MSO_LANGUAGE_ID.ENGLISH_US

    # Test setting to None (python-pptx maps None to MSO_LANGUAGE_ID.NONE)
    font.language_id = None
    assert font.language_id == MSO_LANGUAGE_ID.NONE

    pptx_path = output / "font_language_id.pptx"
    presentation.save(pptx_path)
