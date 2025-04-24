"""Test pptx wrapper module."""

import pytest
from pptx import Presentation as PptxPresentation

from pptxr._pptx import (
    Presentation as PptxWrapper,
)
from pptxr._pptx import (
    Shape as ShapeWrapper,
)
from pptxr._pptx import (
    Slide as SlideWrapper,
)


def test_pptx_wrapper_init():
    """Test PptxWrapper initialization."""
    wrapper = PptxWrapper()
    pptx_obj = wrapper.to_pptx()
    assert hasattr(pptx_obj, "slides")
    assert hasattr(pptx_obj, "slide_layouts")


def test_pptx_wrapper_from_pptx():
    """Test PptxWrapper.from_pptx."""
    pptx_obj = PptxPresentation()
    wrapper = PptxWrapper.from_pptx(pptx_obj)
    assert isinstance(wrapper, PptxWrapper)
    assert wrapper.to_pptx() is pptx_obj

    with pytest.raises(TypeError):
        PptxWrapper.from_pptx("not a presentation")


def test_slide_wrapper_init(tmp_path):
    """Test SlideWrapper initialization."""
    pptx = PptxPresentation()
    pptx.save(tmp_path / "test.pptx")
    pptx = PptxPresentation(tmp_path / "test.pptx")
    slide = pptx.slides.add_slide(pptx.slide_layouts[0])
    wrapper = SlideWrapper(slide)
    assert wrapper.to_pptx() is slide


def test_slide_wrapper_from_pptx(tmp_path):
    """Test SlideWrapper.from_pptx."""
    pptx = PptxPresentation()
    pptx.save(tmp_path / "test.pptx")
    pptx = PptxPresentation(tmp_path / "test.pptx")
    slide = pptx.slides.add_slide(pptx.slide_layouts[0])
    wrapper = SlideWrapper.from_pptx(slide)
    assert isinstance(wrapper, SlideWrapper)
    assert wrapper.to_pptx() is slide

    with pytest.raises(TypeError):
        SlideWrapper.from_pptx("not a slide")


def test_shape_wrapper_init(tmp_path):
    """Test ShapeWrapper initialization."""
    pptx = PptxPresentation()
    pptx.save(tmp_path / "test.pptx")
    pptx = PptxPresentation(tmp_path / "test.pptx")
    slide = pptx.slides.add_slide(pptx.slide_layouts[0])
    shape = slide.shapes.title
    assert shape
    wrapper = ShapeWrapper(shape)
    assert wrapper.to_pptx() is shape


def test_shape_wrapper_from_pptx(tmp_path):
    """Test ShapeWrapper.from_pptx."""
    pptx = PptxPresentation()
    pptx.save(tmp_path / "test.pptx")
    pptx = PptxPresentation(tmp_path / "test.pptx")
    slide = pptx.slides.add_slide(pptx.slide_layouts[0])
    shape = slide.shapes.title
    wrapper = ShapeWrapper.from_pptx(shape)
    assert isinstance(wrapper, ShapeWrapper)
    assert wrapper.to_pptx() is shape

    with pytest.raises(TypeError):
        ShapeWrapper.from_pptx("not a shape")
