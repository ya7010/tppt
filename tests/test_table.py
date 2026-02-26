"""Tests for table module."""

import pathlib
from dataclasses import dataclass

import pytest

import tppt
import tppt.pptx.table
from tppt._features import (
    USE_PANDAS,
    USE_POLARS,
    USE_PYDANTIC,
    Dataclass,
    PydanticModel,
)
from tppt.pptx.table.table import Column, ColumnCollection, Row, RowCollection


def test_create_table_with_list_data(output: pathlib.Path) -> None:
    """Test creating a table with list data."""
    table_data = [
        ["Header 1", "Header 2", "Header 3"],
        ["Cell 1,1", "Cell 1,2", "Cell 1,3"],
        ["Cell 2,1", "Cell 2,2", "Cell 2,3"],
    ]

    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .table(
                table_data,
                left=(100, "pt"),
                top=(100, "pt"),
                width=(400, "pt"),
                height=(200, "pt"),
            )
        )
        .build()
    )
    presentation.save(output / "table_list_data.pptx")


@pytest.mark.skipif(not USE_PANDAS, reason="Pandas not installed")
def test_create_table_with_pandas_dataframe(output: pathlib.Path) -> None:
    """Test creating a table with pandas DataFrame."""
    # Type checking is done at runtime when pandas is available
    import pandas as pd  # type: ignore[import]

    # Create pandas DataFrame with simple data using dictionary
    data = {
        "名前": ["田中", "佐藤", "鈴木"],
        "年齢": [25, 30, 22],
        "都市": ["東京", "大阪", "名古屋"],
    }
    df = pd.DataFrame(data)

    # Convert DataFrame to a list of lists for table creation
    table_data = [df.columns.tolist()] + df.values.tolist()

    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .table(
                table_data,  # 変換したリストを使用
                left=(100, "pt"),
                top=(100, "pt"),
                width=(400, "pt"),
                height=(200, "pt"),
                first_row_header=True,
            )
        )
        .build()
    )
    presentation.save(output / "table_pandas_data.pptx")


@pytest.mark.skipif(not USE_POLARS, reason="Polars not installed")
def test_create_table_with_polars_dataframe(output: pathlib.Path) -> None:
    """Test creating a table with polars DataFrame."""
    # Type checking is done at runtime when polars is available
    import polars as pl  # type: ignore[import]

    # Create polars DataFrame with simple data using dictionary
    data = {
        "製品": ["A製品", "B製品", "C製品"],
        "価格": [1000, 2000, 3000],
        "在庫": [50, 30, 10],
    }
    df = pl.DataFrame(data)

    # Convert DataFrame to a list of lists for table creation
    table_data = [df.columns] + df.to_numpy().tolist()

    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .table(
                table_data,  # 変換したリストを使用
                left=(100, "pt"),
                top=(100, "pt"),
                width=(400, "pt"),
                height=(200, "pt"),
                first_row_header=True,
            )
        )
        .build()
    )
    presentation.save(output / "table_polars_data.pptx")


@pytest.mark.skipif(not USE_POLARS, reason="Polars not installed")
def test_create_table_with_polars_lazyframe(output: pathlib.Path) -> None:
    """Test creating a table with polars LazyFrame."""
    # Type checking is done at runtime when polars is available
    import polars as pl  # type: ignore[import]

    # Create polars LazyFrame with simple data using dictionary
    data = {
        "カテゴリ": ["食品", "電化製品", "衣類"],
        "売上": [5000, 12000, 8000],
        "利益率": [0.2, 0.15, 0.25],
    }
    lazy_df = pl.LazyFrame(data)

    # Collect LazyFrame to DataFrame
    df = lazy_df.collect()

    # Convert DataFrame to a list of lists for table creation
    table_data = [df.columns] + df.to_numpy().tolist()

    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .table(
                table_data,
                left=(100, "pt"),
                top=(100, "pt"),
                width=(400, "pt"),
                height=(200, "pt"),
                first_row_header=True,
            )
        )
        .build()
    )
    presentation.save(output / "table_polars_lazyframe_data.pptx")


def test_create_table_with_dataclass(output: pathlib.Path) -> None:
    """Test creating a table with dataclass."""

    @dataclass
    class Employee:
        """Employee data class for testing."""

        name: str
        department: str
        salary: int
        join_date: str

    # Create employee instances
    employees: list[Dataclass] = [
        Employee("山田太郎", "営業部", 500000, "2020-04-01"),
        Employee("鈴木花子", "開発部", 600000, "2019-07-15"),
        Employee("佐藤一郎", "人事部", 550000, "2021-01-10"),
    ]

    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .table(
                employees,
                left=(100, "pt"),
                top=(100, "pt"),
                width=(400, "pt"),
                height=(200, "pt"),
                first_row_header=True,
            )
        )
        .build()
    )
    presentation.save(output / "table_dataclass_data.pptx")


@pytest.mark.skipif(not USE_PYDANTIC, reason="Pydantic not installed")
def test_create_table_with_pydantic_model(output: pathlib.Path) -> None:
    """Test creating a table with Pydantic model."""

    class Product(PydanticModel):
        """Product model for testing."""

        name: str
        category: str
        price: int
        stock: int

    # Create product instances
    products: list[PydanticModel] = [
        Product(name="ノートPC", category="電子機器", price=150000, stock=10),
        Product(name="スマートフォン", category="電子機器", price=80000, stock=20),
        Product(name="タブレット", category="電子機器", price=50000, stock=15),
    ]

    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .table(
                products,
                left=(100, "pt"),
                top=(100, "pt"),
                width=(400, "pt"),
                height=(200, "pt"),
                first_row_header=True,
            )
        )
        .build()
    )
    presentation.save(output / "table_pydantic_data.pptx")


def test_cell_set_chain_methods(output) -> None:
    """Test Cell set_* chain methods return self."""
    from pptx.enum.text import MSO_VERTICAL_ANCHOR

    from tppt.pptx.table.cell import Cell
    from tppt.types import Points

    def customize_table(
        table: tppt.pptx.table.Table,
    ) -> tppt.pptx.table.Table:
        cell = table.cell(0, 0)

        # Test chaining
        result = (
            cell.set_text("chained")
            .set_margin_left(Points(10))
            .set_margin_right(Points(10))
            .set_margin_top(Points(5))
            .set_margin_bottom(Points(5))
            .set_vertical_anchor(MSO_VERTICAL_ANCHOR.MIDDLE)
        )

        assert isinstance(result, Cell)
        assert result is cell
        assert cell.text == "chained"
        assert cell.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE

        return table

    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .table(
                tppt.apply(customize_table),
                rows=2,
                cols=2,
                left=(100, "pt"),
                top=(100, "pt"),
                width=(400, "pt"),
                height=(200, "pt"),
            )
        )
        .build()
    )
    presentation.save(output / "cell_chain.pptx")


def test_table_rows_columns_banding(output) -> None:
    """Test Table rows, columns, banding properties, and iter_cells."""
    table_data = [
        ["A", "B", "C"],
        ["1", "2", "3"],
        ["4", "5", "6"],
    ]

    def customize_table(
        table: tppt.pptx.table.Table,
    ) -> tppt.pptx.table.Table:
        # Test rows
        rows = table.rows
        assert isinstance(rows, RowCollection)
        assert len(rows) == 3
        row = rows[0]
        assert isinstance(row, Row)
        cells = row.cells
        assert len(cells) == 3

        # Test columns
        columns = table.columns
        assert isinstance(columns, ColumnCollection)
        assert len(columns) == 3
        col = columns[0]
        assert isinstance(col, Column)

        # Test banding properties with chaining
        result = (
            table.set_first_row(True)
            .set_last_row(False)
            .set_first_col(True)
            .set_last_col(False)
            .set_horz_banding(True)
            .set_vert_banding(False)
        )
        assert result is table
        assert table.first_row is True
        assert table.last_row is False
        assert table.first_col is True
        assert table.last_col is False
        assert table.horz_banding is True
        assert table.vert_banding is False

        # Test iter_cells
        cell_count = sum(1 for _ in table.iter_cells())
        assert cell_count == 9  # 3x3

        return table

    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .table(
                table_data,
                left=(100, "pt"),
                top=(100, "pt"),
                width=(400, "pt"),
                height=(200, "pt"),
            )
            .tap(lambda slide: customize_table(
                tppt.pptx.table.Table(slide.to_pptx().slides[0].shapes[0].table)
                if hasattr(slide.to_pptx(), "slides")
                else None
            ) if False else None)
        )
        .build()
    )

    # Access table from the built presentation and test
    pptx_pres = presentation.to_pptx()
    pptx_table = pptx_pres.slides[0].shapes[0].table
    table = tppt.pptx.table.Table(pptx_table)

    rows = table.rows
    assert isinstance(rows, RowCollection)
    assert len(rows) == 3

    columns = table.columns
    assert isinstance(columns, ColumnCollection)
    assert len(columns) == 3

    table.set_first_row(True).set_horz_banding(True)
    assert table.first_row is True
    assert table.horz_banding is True

    cell_count = sum(1 for _ in table.iter_cells())
    assert cell_count == 9

    presentation.save(output / "table_rows_cols.pptx")


def test_table_callback(output) -> None:
    """Test creating table callback override."""

    def customize_table(
        table: tppt.pptx.table.Table,
    ) -> tppt.pptx.table.Table:
        # 3 x 3 のデータとヘッダを作成
        table.cell(0, 0).text = "Price"
        table.cell(0, 1).text = "Quantity"
        table.cell(0, 2).text = "Total"

        table.cell(1, 0).text = "100"
        table.cell(1, 1).text = "5"
        table.cell(1, 2).text = "500"

        table.cell(2, 0).text = "200"
        table.cell(2, 1).text = "10"
        table.cell(2, 2).text = "2000"

        table.cell(3, 0).text = "300"
        table.cell(3, 1).text = "15"
        table.cell(3, 2).text = "4500"

        return table

    # プレゼンテーションの作成
    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .table(
                tppt.apply(customize_table),
                rows=4,
                cols=3,
                left=(100, "pt"),
                top=(100, "pt"),
                width=(400, "pt"),
                height=(200, "pt"),
            )
        )
        .build()
    )

    # プレゼンテーションを保存
    pptx_path = output / "table_callback.pptx"
    presentation.save(pptx_path)
