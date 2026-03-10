"""Tests for presentation module."""

import pathlib
import subprocess
import sys

import pytest

import tppt
from tppt.pptx.table.table import TableCellStyle


def test_create_presentation(output: pathlib.Path) -> None:
    """Test creating a presentation."""
    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .text(
                "text",
                left=(100, "pt"),
                top=(100, "pt"),
                width=(100, "pt"),
                height=(100, "pt"),
            )
        )
        .build()
    )
    presentation.save(output / "text_only.pptx")


def test_presentation_width_and_height(output: pathlib.Path) -> None:
    """Test setting the width and height of a presentation."""
    presentation = tppt.Presentation.builder().build()
    assert presentation.slide_width == tppt.types.EnglishMetricUnits(9144000)
    assert presentation.slide_height == tppt.types.EnglishMetricUnits(6858000)


def test_presentation_width_and_height_setter(output: pathlib.Path) -> None:
    """Test setting the width and height of a presentation."""
    presentation = tppt.Presentation.builder().build()
    presentation.slide_width = (20, "cm")
    presentation.slide_height = (10, "cm")
    assert presentation.slide_width == (20, "cm")
    assert presentation.slide_height == (10, "cm")


def test_create_presentation_with_table(output: pathlib.Path) -> None:
    """Test creating a presentation with a table."""
    table_data = [
        ["Header 1", "Header 2", "Header 3"],
        ["Cell 1,1", "Cell 1,2", "Cell 1,3"],
        ["Cell 2,1", "Cell 2,2", "Cell 2,3"],
    ]

    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .table(
                table_data,
                left=(100, "pt"),
                top=(100, "pt"),
                width=(400, "pt"),
                height=(200, "pt"),
            )
        )
        .build()
    )
    presentation.save(output / "table_only.pptx")


def test_create_presentation_with_styled_table(output: pathlib.Path) -> None:
    """Test creating a presentation with a styled table."""
    data = [
        ["Header 1", "Header 2", "Header 3"],
        ["Cell 1,1", "Cell 1,2", "Cell 1,3"],
        ["Cell 2,1", "Cell 2,2", "Cell 2,3"],
    ]

    cell_styles: list[list[TableCellStyle]] = [
        [
            {"bold": True, "font_size": (14, "pt"), "text_align": "center"},
            {"bold": True, "font_size": (14, "pt"), "text_align": "center"},
            {"bold": True, "font_size": (14, "pt"), "text_align": "center"},
        ],
        [
            {"text_align": "left", "vertical_align": "middle"},
            {"text_align": "center", "vertical_align": "middle", "italic": True},
            {"text_align": "right", "vertical_align": "middle"},
        ],
        [
            {"text_align": "left", "vertical_align": "bottom"},
            {"text_align": "center", "vertical_align": "bottom"},
            {"text_align": "right", "vertical_align": "bottom", "font_name": "Arial"},
        ],
    ]

    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .table(
                data,
                left=(100, "pt"),
                top=(100, "pt"),
                width=(400, "pt"),
                height=(200, "pt"),
                cell_styles=cell_styles,
                first_row_header=True,
            )
        )
        .build()
    )
    presentation.save(output / "styled_table.pptx")


@pytest.mark.parametrize(
    "example_file",
    [f for f in pathlib.Path(__file__).parent.parent.joinpath("examples").glob("*.py")],
)
def test_example_files(example_file: pathlib.Path) -> None:
    """Test that example Python files run without errors.

    Args:
        example_file: Path to the example Python file.
    """
    if example_file.with_suffix(".pptx").exists():
        example_file.with_suffix(".pptx").unlink()

    # Run the example script as a subprocess
    result = subprocess.run(
        [sys.executable, str(example_file)],
        capture_output=True,
        text=True,
        check=False,
    )

    # Check that the script ran successfully
    assert result.returncode == 0, (
        f"Example {example_file.name} failed with error: {result.stderr}"
    )

    # Also check that the output PPTX file was created
    expected_pptx = example_file.with_suffix(".pptx")
    assert expected_pptx.exists(), (
        f"Expected output file {expected_pptx} was not created"
    )


def test_picture_crop_properties(output: pathlib.Path) -> None:
    """Test Picture crop properties and chaining."""
    import io

    from PIL import Image

    from tppt.pptx.shape.picture import Picture

    # Create a minimal test image in memory
    img = Image.new("RGB", (100, 100), color="red")
    img_bytes = io.BytesIO()
    img.save(img_bytes, format="PNG")
    img_bytes.seek(0)

    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .picture(
                img_bytes,
                left=(100, "pt"),
                top=(100, "pt"),
                width=(200, "pt"),
                height=(200, "pt"),
            )
        )
        .build()
    )

    # Access the picture shape
    pptx_slide = presentation.to_pptx().slides[0]
    from pptx.shapes.picture import Picture as PptxPicture

    pptx_shape = pptx_slide.shapes[0]
    picture = Picture(pptx_shape)

    # Test crop setters
    picture.crop_bottom = 0.1
    picture.crop_left = 0.2
    picture.crop_right = 0.15
    picture.crop_top = 0.05

    assert abs(picture.crop_bottom - 0.1) < 0.001
    assert abs(picture.crop_left - 0.2) < 0.001
    assert abs(picture.crop_right - 0.15) < 0.001
    assert abs(picture.crop_top - 0.05) < 0.001

    # Test chaining
    result = (
        picture.set_crop_bottom(0.0)
        .set_crop_left(0.0)
        .set_crop_right(0.0)
        .set_crop_top(0.0)
    )
    assert result is picture

    # Test image property
    image = picture.image
    assert image is not None

    # Test line property
    line = picture.line
    assert line is not None

    presentation.save(output / "picture_crop.pptx")


def test_chart_properties(output: pathlib.Path) -> None:
    """Test Chart properties (chart_style, chart_type, has_legend, has_title, etc.)."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    from tppt.pptx.chart.chart import Chart, ChartTitle, Legend

    chart_data = CategoryChartData()
    chart_data.categories = ["East", "West", "Midwest"]
    chart_data.add_series("Sales", (19.2, 21.4, 16.7))

    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .chart(
                chart_type="Clustered Column",
                x=(50, "pt"),
                y=(50, "pt"),
                cx=(400, "pt"),
                cy=(300, "pt"),
                chart_data=chart_data,
            )
        )
        .build()
    )

    # Access the chart
    pptx_slide = presentation.to_pptx().slides[0]
    pptx_chart = pptx_slide.shapes[0].chart
    chart = Chart(pptx_chart)

    # Test chart_type
    assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED

    # Test chart_style
    chart.chart_style = 2
    assert chart.chart_style == 2
    result = chart.set_chart_style(10)
    assert result is chart
    assert chart.chart_style == 10

    # Test has_legend
    chart.has_legend = True
    assert chart.has_legend is True
    result = chart.set_has_legend(False)
    assert result is chart
    assert chart.has_legend is False

    # Test has_title
    chart.has_title = True
    assert chart.has_title is True
    result = chart.set_has_title(False)
    assert result is chart

    # Test chart_title
    chart.has_title = True
    ct = chart.chart_title
    assert isinstance(ct, ChartTitle)
    assert isinstance(ct.has_text_frame, bool)

    # Test font
    font = chart.font
    assert font is not None

    # Test legend
    chart.has_legend = True
    legend = chart.legend
    assert isinstance(legend, Legend)
    legend_font = legend.font
    assert legend_font is not None

    # Test replace_data
    new_data = CategoryChartData()
    new_data.categories = ["North", "South"]
    new_data.add_series("Revenue", (30.0, 25.0))
    chart.replace_data(new_data)

    presentation.save(output / "chart_props.pptx")
