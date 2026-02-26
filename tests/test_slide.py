"""Tests for slide module."""

import pytest

import tppt
import tppt.pptx.slide
from tppt.pptx.shape import BaseShape
from tppt.pptx.shape.background import Background
from tppt.pptx.shape.placeholder import LayoutPlaceholder, SlidePlaceholder
from tppt.pptx.slide_layout import SlideLayout as PptxSlideLayout


@pytest.mark.skip(reason="This test is not implemented yet.")
def test_slide_placeholders() -> None:
    """Test getting placeholders from a slide."""
    # Create a presentation with a slide
    presentation = (
        tppt.Presentation.builder().slide(lambda slide: slide.BlankLayout()).build()
    )

    # Get the slide from the presentation's pptx object
    pptx_presentation = presentation.to_pptx()
    if not pptx_presentation.slides:
        pytest.skip("No slides in presentation")

    slide = tppt.pptx.slide.Slide(pptx_presentation.slides[0])

    # Test that we can get placeholders
    placeholders = slide.placeholders

    # Check that placeholders is a list
    assert isinstance(placeholders, list)
    assert len(placeholders) != 0

    # Verify each placeholder is a SlidePlaceholder instance
    for placeholder in placeholders:
        assert isinstance(placeholder, SlidePlaceholder)

        # Check that we can convert it to pptx object
        pptx_placeholder = placeholder.to_pptx()
        assert pptx_placeholder is not None


def test_slide_layout_placeholders() -> None:
    """Test getting placeholders from a slide layout."""
    # Create a presentation
    presentation = tppt.Presentation.builder().build()

    # Get the slide layout from the presentation's pptx object
    pptx_presentation = presentation.to_pptx()
    if not pptx_presentation.slide_layouts:
        pytest.skip("No slide layouts in presentation")

    # Get the first slide layout
    slide_layout = PptxSlideLayout(pptx_presentation.slide_layouts[0])

    # Test that we can get placeholders
    placeholders = slide_layout.placeholders

    # Check that placeholders is a list
    assert isinstance(placeholders, list)
    assert len(placeholders) != 0

    # Verify each placeholder is a LayoutPlaceholder instance
    for placeholder in placeholders:
        assert isinstance(placeholder, LayoutPlaceholder)

        # Check that we can convert it to pptx object
        pptx_placeholder = placeholder.to_pptx()
        assert pptx_placeholder is not None


def test_slide_builder_tap(output) -> None:
    """Test that tap() invokes the callback with a tppt Slide wrapper."""
    callback_called = False

    def my_callback(slide: tppt.pptx.slide.Slide) -> None:
        nonlocal callback_called
        callback_called = True
        assert isinstance(slide, tppt.pptx.slide.Slide)
        assert slide.to_pptx() is not None

    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .text(
                "Before tap",
                left=(100, "pt"),
                top=(100, "pt"),
                width=(200, "pt"),
                height=(50, "pt"),
            )
            .tap(my_callback)
            .text(
                "After tap",
                left=(100, "pt"),
                top=(200, "pt"),
                width=(200, "pt"),
                height=(50, "pt"),
            )
        )
        .build()
    )

    assert callback_called
    assert len(presentation.slides[0].shapes) >= 2
    presentation.save(output / "tap_test.pptx")


def test_slide_builder_tap_with_raw_pptx(output) -> None:
    """Test that tap() can modify the slide via the raw python-pptx API."""
    from pptx.util import Inches

    def add_raw_textbox(slide: tppt.pptx.slide.Slide) -> None:
        pptx_slide = slide.to_pptx()
        txBox = pptx_slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        )
        txBox.text_frame.text = "Added via tap()"

    presentation = (
        tppt.Presentation.builder()
        .slide(lambda slide: slide.BlankLayout().builder().tap(add_raw_textbox))
        .build()
    )

    slide = presentation.slides[0]
    texts = [
        shape.to_pptx().text
        for shape in slide.shapes
        if hasattr(shape.to_pptx(), "text")
    ]
    assert any("Added via tap()" in t for t in texts)
    presentation.save(output / "tap_raw_pptx_test.pptx")


def test_background_fill(output) -> None:
    """Test Background.fill property."""
    from tppt.pptx.dml.fill import FillFormat

    presentation = (
        tppt.Presentation.builder()
        .slide(lambda slide: slide.BlankLayout())
        .build()
    )

    slide = presentation.slides[0]
    bg = slide.background
    assert isinstance(bg, Background)

    fill = bg.fill
    assert isinstance(fill, FillFormat)


def test_notes_slide_properties(output) -> None:
    """Test NotesSlide properties."""
    from tppt.pptx.notes_slide import NotesSlide
    from tppt.pptx.text.text_frame import TextFrame

    presentation = (
        tppt.Presentation.builder()
        .slide(lambda slide: slide.BlankLayout())
        .build()
    )

    # Access the pptx slide directly to create a notes slide
    pptx_slide = presentation.to_pptx().slides[0]
    pptx_notes = pptx_slide.notes_slide  # This creates the notes slide

    notes = NotesSlide(pptx_notes)

    # Test notes_text_frame
    text_frame = notes.notes_text_frame
    assert isinstance(text_frame, TextFrame)

    # Test notes_placeholder
    placeholder = notes.notes_placeholder
    assert isinstance(placeholder, SlidePlaceholder)

    # Test placeholders
    placeholders = notes.placeholders
    assert isinstance(placeholders, list)
    assert len(placeholders) > 0

    # Test shapes
    shapes = notes.shapes
    assert isinstance(shapes, list)
    assert len(shapes) > 0
    for shape in shapes:
        assert isinstance(shape, BaseShape)


def test_baseshape_boolean_properties(output) -> None:
    """Test BaseShape boolean properties: has_chart, has_table, has_text_frame, is_placeholder."""
    presentation = (
        tppt.Presentation.builder()
        .slide(
            lambda slide: slide.BlankLayout()
            .builder()
            .text(
                "hello",
                left=(100, "pt"),
                top=(100, "pt"),
                width=(200, "pt"),
                height=(50, "pt"),
            )
        )
        .build()
    )

    slide = presentation.slides[0]
    shapes = slide.shapes
    assert len(shapes) > 0

    for shape in shapes:
        # All boolean properties should be accessible
        assert isinstance(shape.has_chart, bool)
        assert isinstance(shape.has_table, bool)
        assert isinstance(shape.has_text_frame, bool)
        assert isinstance(shape.is_placeholder, bool)

    # Text box should have text frame but not chart/table
    text_shape = shapes[0]
    assert text_shape.has_text_frame is True
    assert text_shape.has_chart is False
    assert text_shape.has_table is False
